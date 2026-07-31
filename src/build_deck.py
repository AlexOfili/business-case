"""
Rebuild the 8-slide consulting deck (PPTX).

Run:
    python3 build_deck.py

Slide list:
  1. Cover
  2. Situation / Complication / Question / Answer
  3. The question — formula
  4. Recommendation
  5. The numbers — 4 price tiers
  6. The curve (with chart)
  7. Risks (sensitivity chart)
  8. Next steps
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(os.path.dirname(HERE), "consulting_deck.pptx")
CHARTS = os.path.join(os.path.dirname(HERE), "charts")

# Brand colours
NAVY  = RGBColor(0x1F, 0x4E, 0x79)
RED   = RGBColor(0xB3, 0x00, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT  = RGBColor(0xE8, 0xEE, 0xF7)


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=None,
             align=None, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_header(slide, num, title, page, total=8):
    add_rect(slide, 0, 0, 13.33, 0.55, NAVY)
    add_text(slide, 0.4, 0.08, 9, 0.4,
             f"{num:02d} · {title}", size=14, bold=True, color=WHITE)
    add_text(slide, 10.5, 0.08, 2.5, 0.4,
             f"{page} / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, 0.4, 7.18, 12, 0.3,
             "Day 9 · Riverside Malls · Smart Trolley Pilot",
             size=9, italic=True, color=GREY)


# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: Cover ----
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.33, 7.5, NAVY)
    add_text(s, 0.6, 0.6, 12, 0.4,
             "DAY 9  ·  CONSULTING RECOMMENDATION",
             size=14, bold=True, color=WHITE)
    add_text(s, 0.6, 1.4, 12, 1.4,
             "Should Riverside Malls roll out 25 smart trolleys — and at what price?",
             size=30, bold=True, color=WHITE)
    add_text(s, 0.6, 3.0, 12, 0.9,
             "Yes — at £200 per trolley per month.",
             size=40, bold=True, color=RGBColor(0xFF, 0xD7, 0x00))
    add_text(s, 0.6, 4.0, 12, 0.6,
             "£1,000 / month profit   ·   6-trolley margin of safety   ·   break-even at 18.75 trolleys",
             size=18, color=WHITE)
    add_text(s, 0.6, 6.8, 12, 0.3,
             "Knull Work Experience Programme   ·   Confidential to participants",
             size=10, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC))

    # ---- Slide 2: SCQA ----
    s = prs.slides.add_slide(blank)
    add_header(s, 1, "SITUATION", 2)
    add_text(s, 0.4, 0.9, 12.5, 0.6,
             "Riverside Malls is exploring a smart-trolley pilot",
             size=22, bold=True, color=NAVY)
    blocks = [
        ("SITUATION",
         "Riverside Malls wants a tech-forward differentiator and is "
         "considering a fleet of 25 smart trolleys to give shoppers a "
         "frictionless cart and ops a shopper-flow signal."),
        ("COMPLICATION",
         "Hardware + telematics + support is real money. A naive rollout "
         "could lose money every month — or worse, prove the unit economics "
         "don't work before the next budget cycle."),
        ("QUESTION",
         "Should we go ahead with 25 trolleys, and at what price per "
         "trolley per month?"),
        ("ANSWER",
         "Yes — at £200 per trolley per month. £1,000/month profit with a "
         "6-trolley margin of safety above break-even."),
    ]
    y = 1.7
    for label, body in blocks:
        add_rect(s, 0.4, y, 1.6, 1.0, SOFT)
        add_text(s, 0.5, y + 0.25, 1.4, 0.5, label, size=11, bold=True, color=NAVY)
        add_text(s, 2.1, y, 10.8, 1.0, body, size=13, color=GREY)
        y += 1.2

    # ---- Slide 3: The question ----
    s = prs.slides.add_slide(blank)
    add_header(s, 2, "THE QUESTION", 3)
    add_text(s, 0.4, 0.9, 12.5, 0.6,
             "Two numbers decide everything",
             size=22, bold=True, color=NAVY)
    add_text(s, 0.4, 1.7, 12.5, 0.5,
             "What price per trolley per month turns a 25-trolley pilot into a profitable one?",
             size=15, italic=True, color=GREY)
    # Revenue
    add_rect(s, 0.4, 2.4, 6.2, 2.6, WHITE, line=NAVY)
    add_text(s, 0.6, 2.5, 5.8, 0.4, "REVENUE   (per trolley / month)", size=12, bold=True, color=NAVY)
    add_text(s, 0.6, 2.9, 5.8, 0.5, "Operator fee:   P   £", size=18, bold=True, color=RED)
    add_text(s, 0.6, 3.5, 5.8, 1.5,
             "•  Charged per active trolley, billed monthly.\n"
             "•  Only revenue counted in the base case.\n"
             "•  £1/visit shopper hire and ad slots are upside — called out, not booked.",
             size=11, color=GREY)
    # Cost
    add_rect(s, 6.8, 2.4, 6.2, 2.6, WHITE, line=NAVY)
    add_text(s, 7.0, 2.5, 5.8, 0.4, "COST   (monthly)", size=12, bold=True, color=NAVY)
    add_text(s, 7.0, 2.9, 5.8, 0.5, "v = £40   ·   F = £3,000   ·   T = 25",
             size=18, bold=True, color=NAVY)
    add_text(s, 7.0, 3.5, 5.8, 1.5,
             "•  v = variable cost per trolley (maintenance, telematics, support).\n"
             "•  F = fixed monthly cost (platform, depot, on-call).\n"
             "•  T = fleet size in the pilot = 25 trolleys.",
             size=11, color=GREY)
    # Formula
    add_rect(s, 0.4, 5.2, 12.6, 1.6, SOFT)
    add_text(s, 0.6, 5.3, 12.2, 0.4, "THE FORMULA", size=12, bold=True, color=NAVY)
    add_text(s, 0.6, 5.7, 12.2, 0.5,
             "Monthly profit = T × (P − v) − F        Break-even trolleys n* = F / (P − v)",
             size=18, bold=True, color=NAVY)
    add_text(s, 0.6, 6.3, 12.2, 0.4,
             "Change one number, every cell updates.  (See unit_economics.xlsx — Assumptions sheet.)",
             size=11, italic=True, color=GREY)

    # ---- Slide 4: Recommendation ----
    s = prs.slides.add_slide(blank)
    add_header(s, 3, "RECOMMENDATION", 4)
    add_text(s, 0.4, 0.9, 12.5, 0.6,
             "Run the pilot. Charge £200 per trolley per month.",
             size=22, bold=True, color=NAVY)
    # Big number block
    add_rect(s, 0.4, 1.8, 4.0, 2.4, NAVY)
    add_text(s, 0.4, 2.0, 4.0, 0.4, "MONTHLY PROFIT", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 2.5, 4.0, 1.0, "£1,000", size=60, bold=True, color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.CENTER)
    add_text(s, 0.4, 3.6, 4.0, 0.4, "at £200 / trolley / month", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 4.6, 1.8, 4.0, 2.4, NAVY)
    add_text(s, 4.6, 2.0, 4.0, 0.4, "BREAK-EVEN TROLLEYS", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 4.6, 2.5, 4.0, 1.0, "18.75", size=60, bold=True, color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.CENTER)
    add_text(s, 4.6, 3.6, 4.0, 0.4, "vs. 25 in the pilot → 6-trolley safety", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 8.8, 1.8, 4.0, 2.4, NAVY)
    add_text(s, 8.8, 2.0, 4.0, 0.4, "HEADROOM TO NEXT TIER", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 8.8, 2.5, 4.0, 1.0, "+30%", size=60, bold=True, color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.CENTER)
    add_text(s, 8.8, 3.6, 4.0, 0.4, "before £250 is needed", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # Why not £150
    add_rect(s, 0.4, 4.6, 12.6, 2.2, RGBColor(0xFF, 0xF5, 0xF5), line=RED)
    add_text(s, 0.6, 4.7, 12.2, 0.4, "WHY NOT £150?", size=12, bold=True, color=RED)
    add_text(s, 0.6, 5.1, 12.2, 1.5,
             "At £150 the pilot loses £250/month — break-even needs 27.3 trolleys but "
             "we only have 25. Penetration pricing doesn't work until the fixed cost "
             "is amortised across more units.",
             size=14, color=GREY)

    # ---- Slide 5: The numbers ----
    s = prs.slides.add_slide(blank)
    add_header(s, 4, "THE NUMBERS", 5)
    add_text(s, 0.4, 0.9, 12.5, 0.5,
             "Unit economics across 4 price tiers",
             size=22, bold=True, color=NAVY)
    # Table
    rows = [
        ("Price tier", "Price (£)", "Contribution (£)",
         "Monthly profit (£)", "Break-even trolleys"),
        ("Floor (£150)",        "£150", "£110",  "£-250",  "27.3"),
        ("Base (£200)",         "£200", "£160",  "£1,000", "18.75"),
        ("Premium (£250)",      "£250", "£210",  "£2,250", "14.3"),
        ("Aspirational (£300)", "£300", "£260",  "£3,500", "11.5"),
    ]
    col_x  = [0.4, 3.5, 5.6, 7.6, 10.6]
    col_w  = [3.0, 2.0, 2.0, 3.0, 2.4]
    for ri, row in enumerate(rows):
        y = 1.7 + ri * 0.55
        fill = NAVY if ri == 0 else (SOFT if row[0].startswith("Base") else WHITE)
        color = WHITE if ri == 0 else GREY
        for ci, (x, w, val) in enumerate(zip(col_x, col_w, row)):
            add_rect(s, x, y, w, 0.5, fill, line=NAVY if ri == 0 else None)
            add_text(s, x + 0.1, y + 0.1, w - 0.2, 0.3, val,
                     size=13, bold=(ri == 0 or row[0].startswith("Base")),
                     color=color, align=PP_ALIGN.CENTER)
    # Worked example
    add_rect(s, 0.4, 4.7, 12.6, 2.2, SOFT)
    add_text(s, 0.6, 4.8, 12.2, 0.4, "WORKED EXAMPLE  ·  £200", size=12, bold=True, color=NAVY)
    add_text(s, 0.6, 5.2, 12.2, 0.5,
             "contribution = P − v = 200 − 40 = £160",
             size=15, color=GREY)
    add_text(s, 0.6, 5.7, 12.2, 0.5,
             "monthly profit = T × (P − v) − F = 25 × 160 − 3,000 = £1,000",
             size=15, color=GREY)
    add_text(s, 0.6, 6.2, 12.2, 0.5,
             "break-even trolleys n* = F / (P − v) = 3,000 / 160 = 18.75   →   "
             "pilot clears break-even with 6.25 trolleys to spare.",
             size=15, bold=True, color=NAVY)

    # ---- Slide 6: The curve ----
    s = prs.slides.add_slide(blank)
    add_header(s, 5, "THE CURVE", 6)
    add_text(s, 0.4, 0.9, 12.5, 0.5,
             "£200 is the knee of the curve — and the only price with margin",
             size=20, bold=True, color=NAVY)
    s.shapes.add_picture(os.path.join(CHARTS, "profit_vs_price.png"),
                         Inches(0.3), Inches(1.6), Inches(8.0), Inches(5.0))
    add_rect(s, 8.6, 1.6, 4.6, 2.4, SOFT)
    add_text(s, 8.8, 1.7, 4.2, 0.4, "READ", size=12, bold=True, color=NAVY)
    add_text(s, 8.8, 2.1, 4.2, 1.8,
             "The red dot is £200 — the first price at which the pilot turns "
             "green. Below it, every month is a loss. Above it, every extra "
             "£50/trolley adds another £1,250/month in profit.",
             size=11, color=GREY)
    add_rect(s, 8.6, 4.2, 4.6, 2.4, RGBColor(0xFF, 0xF5, 0xF5), line=RED)
    add_text(s, 8.8, 4.3, 4.2, 0.4, "GUARDRAIL", size=12, bold=True, color=RED)
    add_text(s, 8.8, 4.7, 4.2, 1.8,
             "If price falls to £160 or below, profit goes negative. Anchor the "
             "sales conversation at £200 — don't negotiate on the headline "
             "number without a fleet-size concession.",
             size=11, color=GREY)

    # ---- Slide 7: Risks ----
    s = prs.slides.add_slide(blank)
    add_header(s, 6, "RISKS", 7)
    add_text(s, 0.4, 0.9, 12.5, 0.5,
             "The pilot survives 2× support cost — but not 20% monthly churn",
             size=18, bold=True, color=NAVY)
    s.shapes.add_picture(os.path.join(CHARTS, "sensitivity.png"),
                         Inches(0.3), Inches(1.6), Inches(8.0), Inches(5.0))
    add_rect(s, 8.6, 1.6, 4.6, 1.8, SOFT)
    add_text(s, 8.8, 1.7, 4.2, 0.4, "Support cost doubles", size=12, bold=True, color=NAVY)
    add_text(s, 8.8, 2.1, 4.2, 1.2,
             "We drop to exactly £0 profit. We must monitor cost-to-serve weekly "
             "in the first 90 days.",
             size=11, color=GREY)
    add_rect(s, 8.6, 3.6, 4.6, 1.8, RGBColor(0xFF, 0xF5, 0xF5), line=RED)
    add_text(s, 8.8, 3.7, 4.2, 0.4, "20% monthly trolley churn", size=12, bold=True, color=RED)
    add_text(s, 8.8, 4.1, 4.2, 1.2,
             "Lose £400/month. A re-deployment + refurb capex plan is non-negotiable "
             "before we go live.",
             size=11, color=GREY)
    add_rect(s, 8.6, 5.6, 4.6, 1.0, WHITE, line=NAVY)
    add_text(s, 8.8, 5.7, 4.2, 0.8,
             "Trolley under-utilisation: if <18.75 trolleys are active, the pilot "
             "loses money. Recovery playbook + operator co-marketing required.",
             size=10, color=GREY)

    # ---- Slide 8: Next steps ----
    s = prs.slides.add_slide(blank)
    add_header(s, 7, "NEXT STEPS", 8)
    add_text(s, 0.4, 0.9, 12.5, 0.5,
             "Three things to validate before we sign the contract",
             size=22, bold=True, color=NAVY)
    items = [
        ("1", "Confirm cost-to-serve",
         "Shadow-quote support and telematics for 4 weeks with a real depot. The £40 "
         "variable cost is an estimate — moving it to £60 changes break-even from 19 to 25 trolleys."),
        ("2", "Pilot a 5-trolley shadow fleet",
         "Run a no-charge 4-week shadow deployment to measure actual utilisation and churn. "
         "Use the real data, not the model assumption of 100% utilisation."),
        ("3", "Lock the £200 headline price",
         "Set the operator-fee schedule at £200/trolley/month with a contractual volume "
         "clause: any sub-19-trolley month triggers a price-floor review."),
    ]
    y = 1.7
    for num, title, body in items:
        add_rect(s, 0.4, y, 1.0, 1.5, NAVY)
        add_text(s, 0.4, y + 0.3, 1.0, 0.8, num, size=44, bold=True,
                 color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.CENTER)
        add_text(s, 1.6, y + 0.05, 11.4, 0.5, title, size=18, bold=True, color=NAVY)
        add_text(s, 1.6, y + 0.55, 11.4, 1.0, body, size=12, color=GREY)
        y += 1.7

    prs.save(OUT)
    print(f"  wrote {OUT}")
    return OUT


if __name__ == "__main__":
    print("Rebuilding consulting_deck.pptx...")
    build()
    print("Done.")
